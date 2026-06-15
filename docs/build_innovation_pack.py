"""
Generate the ATTEST GCID Innovation Engine pitch pack by cloning the provided
template and replacing its text with ATTEST content. All styling, layouts,
images and graphics from the template are preserved.

Run:
    python docs/build_innovation_pack.py
"""
import copy
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ATTEST diagram palette (matches template teal)
TEAL = RGBColor(0x1C, 0x5E, 0x73)
TEAL_DARK = RGBColor(0x14, 0x47, 0x57)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0x4A, 0x90, 0xA6)

TEMPLATE = Path(r"c:\Users\mbathinedi\Downloads\GCID Innovation Engine - Modernization Pack.pptx")
OUT_DIR = Path(__file__).resolve().parent
OUT = OUT_DIR / "ATTEST_Innovation_Pack.pptx"


# ---------------------------------------------------------------------------
# Content for ATTEST
# ---------------------------------------------------------------------------
TITLE = "ATTEST \u2013 Agent Testing & Trust Evaluation Suite"
AUTHORS = "Manikanta Bathinedi"
DATE = "Date: June 12, 2026"

SUMMARY = (
    "A prompt-driven, framework-agnostic testing and trust-evaluation suite for AI "
    "agents \u2014 standardizing how teams validate agent behavior, tool usage, routing, "
    "and response quality through deterministic assertions, LLM-based evaluators, and "
    "golden baselines \u2014 delivering consistent, secure, and CI-ready quality gates "
    "across agent projects."
)

PROBLEM = (
    "AI agents are shipped today with little to no systematic testing. Teams validate "
    "agents manually and inconsistently \u2014 there is no standard way to verify tool calls, "
    "routing decisions, response quality, safety, or regressions. This leads to "
    "unpredictable agent behavior, silent regressions, and production incidents."
)

WHY = (
    "As AI agents move into production, trust and reliability become critical. ATTEST "
    "brings software-engineering rigor to agent quality \u2014 combining deterministic checks "
    "with LLM-judged evaluations and golden baselines \u2014 so teams catch regressions early, "
    "enforce quality gates in CI/CD, and ship agents with confidence."
)

FOR_WHOM = (
    "AI / Agent Development teams building with Microsoft Agent Framework, Foundry, GHCP / Copilot agents\n"
    "QA & Platform Engineers establishing automated agent quality gates\n"
    "Architects & Tech Leads governing agent reliability, safety, and trust"
)

# Slide 4 metric tiles
M1_PCT = "~30-40%"
M1_SUB = "Manual agent testing & regression triage"
M2_PCT = "~20-25%"
M2_SUB = "Manual QA reviews & escaped defects"
M3_PCT = "~15-20%"
M3_SUB = "Production incidents & LLM eval spend (caching)"

SOLUTION = (
    "ATTEST is an end-to-end agent testing framework that standardizes how agents are "
    "validated before and after release.\n"
    "32 deterministic assertions verify tool calls, JSON, routing, content, PII, cost & performance\n"
    "36 LLM-based evaluators across 4 backends (built-in, DeepEval, Azure AI & RAGAS) score quality, safety & groundedness\n"
    "Multi-turn conversations, user simulation & security red-teaming (30 attacks) stress-test agents\n"
    "Golden baselines detect regressions; web dashboard + CLI + CI deliver repeatable quality gates"
)

BUSINESS_VALUE = (
    "Improves agent reliability by catching regressions before production\n"
    "Enforces consistent, repeatable quality gates across teams\n"
    "Reduces risk and safety gaps through automated evaluation\n"
    "Accelerates delivery with CI-ready, automated agent testing\n"
    "Provides auditable, measurable trust scores for every agent"
)

KEY_TECH = (
    "Deterministic Assertion Engine \u2013 32 checks for tool calls, JSON, routing, content, PII, cost & performance\n"
    "LLM Evaluator Suite \u2013 36 evaluators across 4 backends: built-in, DeepEval, Azure AI Evaluation & RAGAS\n"
    "9 Pluggable Adapters \u2013 Foundry, HTTP REST, Callable, LangChain, LangGraph, CrewAI, AutoGen, OpenAI Assistants & MCP (plus offline mock)\n"
    "Multi-turn, user simulation, security red-teaming (30 attacks) & multi-agent routing assertions\n"
    "Baseline / golden regression + Enterprise Auth (5-tier) + FastAPI dashboard + Typer CLI + CI/CD integration"
)

STATE = (
    "Core framework, 32 assertions, 36 evaluators (4 backends), 9 adapters, security red-teaming, baselines & 9-page dashboard are built and working\n"
    "Next Steps\n"
    "Pilot on real production agents\n"
    "Broaden multi-agent / A2A & MCP scenario coverage\n"
    "Integrate quality gates into CI/CD pipelines"
)

RISKS = (
    "Risks\n"
    "LLM Evaluator Variability\n"
    "LLM-judged scores can vary across runs; mitigated with deterministic assertions and golden baselines\n"
    "Evaluation Cost\n"
    "Large suites can incur LLM cost; mitigated via response caching, rate limiting & cost ceilings\n"
    "\n"
    "Key Dependencies\n"
    "Model Access (Azure OpenAI / Foundry)\n"
    "Evaluators require model endpoints to score agent responses\n"
    "Agent Accessibility\n"
    "Agents must be reachable via a supported adapter (Foundry / HTTP / Callable / framework adapters)\n"
    "GHCP / Copilot License\n"
    "Required for Copilot agent integration scenarios"
)

DEMO_FILE = "ATTEST_Dashboard_Demo.mp4  (recording to be added)"

ASKS = (
    "Effort Required for Prototype:\n"
    "Effort primarily for:\n"
    "Expanding adapter & evaluator coverage\n"
    "Building sample agent test suites\n"
    "Integrating with CI/CD and Foundry\n"
    "\n"
    "Estimated Cost to maintain the Prototype:\n"
    "Minimal \u2014 open-source Python stack; main cost is Azure OpenAI / Foundry usage for "
    "evaluators (controlled via caching & rate limits)\n"
    "\n"
    "Any subscription requirement:  Azure OpenAI / Foundry access; GHCP / Copilot license"
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def iter_text_frames(shapes):
    """Yield every text_frame in shapes, recursing into groups."""
    for shape in shapes:
        if shape.shape_type == 6:  # GROUP
            yield from iter_text_frames(shape.shapes)
        elif shape.has_text_frame:
            yield shape.text_frame


def set_text_preserve(tf, new_text):
    """Replace a text frame's content with new_text (\\n => paragraphs),
    preserving the font formatting of the first run."""
    paras = tf.paragraphs
    first_para = paras[0]
    # Capture style from the first run if present
    style = {}
    if first_para.runs:
        r0 = first_para.runs[0]
        f = r0.font
        style = {
            "size": f.size,
            "bold": f.bold,
            "italic": f.italic,
            "underline": f.underline,
            "name": f.name,
        }
        try:
            if f.color and f.color.type is not None:
                style["rgb"] = f.color.rgb
        except Exception:
            pass

    lines = new_text.split("\n")

    # Remove all paragraphs except the first (operate on XML)
    p_elements = [p._p for p in tf.paragraphs]
    for p in p_elements[1:]:
        p.getparent().remove(p)

    def apply_style(run):
        if style.get("size") is not None:
            run.font.size = style["size"]
        if style.get("bold") is not None:
            run.font.bold = style["bold"]
        if style.get("italic") is not None:
            run.font.italic = style["italic"]
        if style.get("underline") is not None:
            run.font.underline = style["underline"]
        if style.get("name"):
            run.font.name = style["name"]
        if style.get("rgb") is not None:
            run.font.color.rgb = style["rgb"]

    # First paragraph: clear runs, add one
    fp = tf.paragraphs[0]
    for r in list(fp.runs):
        r._r.getparent().remove(r._r)
    run = fp.add_run()
    run.text = lines[0]
    apply_style(run)

    # Remaining lines as new paragraphs
    for line in lines[1:]:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        apply_style(run)


def replace_by_prefix(prs, mapping):
    """For every text frame, if stripped text starts with a key in mapping,
    replace with the mapped value. Keys may match multiple shapes (e.g. the
    title appears on more than one slide)."""
    matched = set()
    for slide in prs.slides:
        for tf in iter_text_frames(slide.shapes):
            txt = tf.text.strip()
            for key, val in mapping.items():
                if txt.startswith(key):
                    set_text_preserve(tf, val)
                    matched.add(key)
                    break
    missing = [k for k in mapping if k not in matched]
    if missing:
        print("WARNING: these keys were not matched:")
        for m in missing:
            print("   -", repr(m[:60]))


# ---------------------------------------------------------------------------
# ATTEST architecture diagram (native shapes, replaces template Picture 6)
# ---------------------------------------------------------------------------
def add_box(slide, l, t, w, h, title, body, title_size=12, body_size=8.5,
            fill=TEAL):
    """Add a rounded teal box with a bold title and bullet body text."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = LINE
    box.line.width = Pt(1)
    box.shadow.inherit = False

    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Emu(91440)
    tf.margin_right = Emu(91440)
    tf.margin_top = Emu(45720)
    tf.margin_bottom = Emu(45720)

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.bold = True
    r.font.size = Pt(title_size)
    r.font.color.rgb = WHITE
    r.font.name = "Segoe UI"

    for line in body:
        bp = tf.add_paragraph()
        bp.alignment = PP_ALIGN.LEFT
        br = bp.add_run()
        br.text = "\u2022 " + line
        br.font.size = Pt(body_size)
        br.font.color.rgb = WHITE
        br.font.name = "Segoe UI"
    return box


def add_arrow(slide, x1, y1, x2, y2):
    """Add a straight connector with an arrowhead at the end."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = LINE
    conn.line.width = Pt(2)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return conn


def build_attest_diagram(slide, L, T, W, H):
    """Draw a vertical ATTEST testing flowchart inside the given EMU box,
    mirroring the template's top -> two-parallel -> converge -> output flow."""
    def y(frac):
        return Emu(int(T + frac * H))

    def wpx(frac):
        return Emu(int(frac * W))

    def hpx(frac):
        return Emu(int(frac * H))

    cx = L + W / 2

    # 1) Top: Test Framework
    a_w, a_h = wpx(0.74), hpx(0.13)
    a_l = Emu(int(cx - a_w / 2))
    a_t = y(0.00)
    add_box(slide, a_l, a_t, a_w, a_h,
            "ATTEST Test Framework",
            ["YAML test scenarios", "Typer CLI + Web Dashboard"],
            title_size=12, body_size=8)

    # 2) Two parallel evaluation boxes
    b_w, b_h = wpx(0.47), hpx(0.30)
    b_t = y(0.20)
    b_l = Emu(int(L + 0.01 * W))
    c_l = Emu(int(L + W - 0.01 * W - b_w))
    add_box(slide, b_l, b_t, b_w, b_h,
            "Assertion Engine (32)",
            ["Tool Calls", "JSON Schema", "Routing & Handoff",
             "Content · PII · Cost", "Performance", "Baseline / Golden"],
            title_size=11, body_size=8)
    add_box(slide, c_l, b_t, b_w, b_h,
            "LLM Evaluators (36)",
            ["Built-in (5)", "DeepEval (12)", "Azure AI (15)",
             "RAGAS (4)", "Safety & Bias", "Groundedness"],
            title_size=11, body_size=8)

    # 3) Converge: Agent Adapters
    d_w, d_h = wpx(0.78), hpx(0.18)
    d_l = Emu(int(cx - d_w / 2))
    d_t = y(0.58)
    add_box(slide, d_l, d_t, d_w, d_h,
            "Agent Under Test \u2014 9 Adapters",
            ["Foundry · HTTP · Callable · LangChain · LangGraph · CrewAI · AutoGen · OpenAI · MCP",
             "5-tier Enterprise Auth (Key \u2192 SP \u2192 WIF \u2192 MI \u2192 Default)"],
            title_size=11, body_size=8, fill=TEAL_DARK)

    # 4) Output: Quality Gates & Reports
    e_w, e_h = wpx(0.86), hpx(0.17)
    e_l = Emu(int(cx - e_w / 2))
    e_t = y(0.82)
    add_box(slide, e_l, e_t, e_w, e_h,
            "Quality Gates & Reports",
            ["Dashboard \u00b7 HTML / JUnit \u00b7 CI/CD gates",
             "Regression diff \u2192 Trusted, Reliable Agent"],
            title_size=11, body_size=8)

    # Connectors (branch -> merge -> output)
    a_bottom = int(a_t + a_h)
    add_arrow(slide, Emu(int(cx)), Emu(a_bottom),
              Emu(int(b_l + b_w / 2)), b_t)
    add_arrow(slide, Emu(int(cx)), Emu(a_bottom),
              Emu(int(c_l + b_w / 2)), b_t)
    b_bottom = int(b_t + b_h)
    add_arrow(slide, Emu(int(b_l + b_w / 2)), Emu(b_bottom),
              Emu(int(cx)), d_t)
    add_arrow(slide, Emu(int(c_l + b_w / 2)), Emu(b_bottom),
              Emu(int(cx)), d_t)
    add_arrow(slide, Emu(int(cx)), Emu(int(d_t + d_h)),
              Emu(int(cx)), e_t)


def replace_diagram(prs):
    """Find the large diagram image on the Technical Details slide and
    replace it with the native ATTEST flowchart in the same position."""
    slide = prs.slides[4]  # slide 5
    target = None
    for s in slide.shapes:
        if s.shape_type == 13:  # PICTURE
            if target is None or (s.width * s.height) > (target.width * target.height):
                target = s
    if target is None:
        print("WARNING: no picture found on slide 5 to replace")
        return
    L, T, W, H = target.left, target.top, target.width, target.height
    target._element.getparent().remove(target._element)
    build_attest_diagram(slide, L, T, W, H)


def main():
    shutil.copyfile(TEMPLATE, OUT)
    prs = Presentation(OUT)

    mapping = {
        # Slide 1 + 2 title
        "Modernization Control Plane": TITLE,
        "Chandan Kumar & Fenil Doshi": AUTHORS,
        "Date: May 6, 2026": DATE,
        "[A prompt-driven": SUMMARY,
        # Slide 3
        "Modernization efforts today are inconsistent": PROBLEM,
        "This idea addresses recurring issues": WHY,
        "Application Development & Modernization Teams": FOR_WHOM,
        # Slide 4 metrics
        "~20-30%": M1_PCT,
        "Per project / per month": M1_SUB,
        "~15-20%": M2_PCT,
        "Manual reviews & support tickets": M2_SUB,
        "[~10-15%]": M3_PCT,
        "Infrastructure & licensing": M3_SUB,
        "The Modernization Control Plane standardizes": SOLUTION,
        "Improves delivery efficiency": BUSINESS_VALUE,
        # Slide 5
        "Prompt Packs Framework": KEY_TECH,
        "Assessment Pack is ready": STATE,
        # Slide 6
        "Risks": RISKS,
        # Slide 7
        "Modernization_Control_Plane_Framework": DEMO_FILE,
        # Slide 8
        "Effort Required for Prototype:": ASKS,
    }

    replace_by_prefix(prs, mapping)
    replace_diagram(prs)
    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()
